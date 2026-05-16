import streamlit as st
from groq import Groq
import pdfplumber
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io, json, os, base64
from datetime import datetime, timezone, timedelta
from concurrent.futures import ThreadPoolExecutor, as_completed

# ── KST 시간 ──────────────────────────────────────────────────────────────────
KST = timezone(timedelta(hours=9))

def now_kst():
    return datetime.now(KST).strftime("%Y.%m.%d %H:%M")

# ── 페이지 설정 ──────────────────────────────────────────────────────────────
st.set_page_config(page_title="심의의견 분석 도구", page_icon="📋", layout="wide")
st.title("📋 보험 광고 심의의견 분석 도구")

# ── API 키 ───────────────────────────────────────────────────────────────────
def get_api_key():
    try:
        return st.secrets["GROQ_API_KEY"]
    except Exception:
        return None

api_key = get_api_key()
if not api_key:
    with st.sidebar:
        st.header("⚙️ 설정")
        api_key = st.text_input("Groq API Key", type="password", help="gsk_...로 시작하는 키")
        if api_key:
            st.success("API 키 입력 완료!")

if not api_key:
    st.info("👈 왼쪽 사이드바에 Groq API 키를 입력하면 시작할 수 있어요.")
    st.stop()

client = Groq(api_key=api_key)

# ── 히스토리 ─────────────────────────────────────────────────────────────────
HISTORY_FILE = "shimui_history.json"

def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return []
    return []

def save_history(history):
    with open(HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=2)

# ── 매체 설정 ────────────────────────────────────────────────────────────────
MEDIA_FILE = "media_settings.json"

DEFAULT_MEDIA = [
    {"media": "네이버 GFA", "placement": "커뮤니케이션애즈", "category": "광고문구 1", "notes": "윗배너", "limit": 40, "include_spaces": True},
    {"media": "당근", "placement": "배너", "category": "브랜드이름", "notes": "브랜드명+보종명까지 같이 기입 가능", "limit": 20, "include_spaces": True},
    {"media": "당근", "placement": "배너", "category": "광고 제목", "notes": "광고 문구", "limit": 30, "include_spaces": True},
    {"media": "토스", "placement": "배너", "category": "주요 문구", "notes": "광고 문구 (윗줄)", "limit": 18, "include_spaces": True},
    {"media": "토스", "placement": "배너", "category": "보조문구", "notes": "광고 문구 (아랫줄에 작게 들어감)", "limit": 18, "include_spaces": True},
]

def load_media():
    if os.path.exists(MEDIA_FILE):
        try:
            with open(MEDIA_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return DEFAULT_MEDIA.copy()

def save_media(media_list):
    with open(MEDIA_FILE, "w", encoding="utf-8") as f:
        json.dump(media_list, f, ensure_ascii=False, indent=2)

# ── OCR ──────────────────────────────────────────────────────────────────────
def ocr_image(image_bytes: bytes) -> str:
    b64 = base64.b64encode(image_bytes).decode("utf-8")
    try:
        resp = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "이 이미지에 적힌 한국어/영어 텍스트를 그대로 추출해주세요.\n"
                            "규칙:\n"
                            "- 보이는 글자만, 있는 그대로, 절대 추측하지 마세요.\n"
                            "- 텍스트가 없으면 '(텍스트 없음)'만 출력하세요.\n"
                            "- 설명, 해석, 부연 설명 절대 금지."
                        )
                    },
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
                ]
            }],
            temperature=0,
            max_tokens=500
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        return f"(OCR 오류: {e})"

def extract_file_text(uploaded_file, use_ocr: bool) -> str:
    file_bytes = uploaded_file.read()
    file_text = ""

    if uploaded_file.name.lower().endswith(".pdf"):
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                file_text += f"[{i}페이지]\n{text}\n\n"

    elif uploaded_file.name.lower().endswith(".pptx"):
        prs = Presentation(io.BytesIO(file_bytes))
        image_jobs = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"
                if use_ocr and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = shape.image.blob
                    image_jobs.append((slide_idx, img_bytes))

            if slide_text:
                file_text += f"[슬라이드 {slide_idx}]\n{slide_text}\n"

        if use_ocr and image_jobs:
            ocr_results = {}
            with ThreadPoolExecutor(max_workers=6) as executor:
                future_map = {
                    executor.submit(ocr_image, img_bytes): slide_idx
                    for slide_idx, img_bytes in image_jobs
                }
                for future in as_completed(future_map):
                    slide_idx = future_map[future]
                    result = future.result()
                    if result and result != "(텍스트 없음)":
                        ocr_results.setdefault(slide_idx, []).append(result)

            for slide_idx, texts in sorted(ocr_results.items()):
                file_text += f"[슬라이드 {slide_idx} - 이미지 텍스트]\n" + "\n".join(texts) + "\n"

    return file_text

# ── 세션 상태 초기화 ──────────────────────────────────────────────────────────
if "selected_history" not in st.session_state:
    st.session_state["selected_history"] = None
if "selected_media_idx" not in st.session_state:
    st.session_state["selected_media_idx"] = None
if "editing_idx" not in st.session_state:
    st.session_state["editing_idx"] = None
if "renaming_history_id" not in st.session_state:
    st.session_state["renaming_history_id"] = None

# ── 사이드바: 히스토리 ────────────────────────────────────────────────────────
history = load_history()

with st.sidebar:
    st.header("📁 최근 작업 히스토리")
    if not history:
        st.caption("아직 작업 내역이 없어요.")
    else:
        for i, item in enumerate(history, 1):
            job_name  = item.get("job_name", "(이름 없음)")
            filename  = item.get("filename", "없음")
            timestamp = item.get("timestamp", "")
            item_id   = item["id"]

            with st.expander(f"{i}. {job_name}"):
                st.caption(f"🕐 {timestamp}")
                st.caption(f"📎 {filename}")

                if st.session_state["renaming_history_id"] == item_id:
                    new_name = st.text_input(
                        "새 작업명",
                        value=job_name,
                        key=f"rename_input_{item_id}"
                    )
                    col_save, col_cancel = st.columns(2)
                    with col_save:
                        if st.button("저장", key=f"rename_save_{item_id}", use_container_width=True):
                            for h in history:
                                if h["id"] == item_id:
                                    h["job_name"] = new_name.strip() or job_name
                            save_history(history)
                            st.session_state["renaming_history_id"] = None
                            st.rerun()
                    with col_cancel:
                        if st.button("취소", key=f"rename_cancel_{item_id}", use_container_width=True):
                            st.session_state["renaming_history_id"] = None
                            st.rerun()
                else:
                    if st.button("이 작업 불러오기", key=f"load_{item_id}", use_container_width=True):
                        st.session_state["selected_history"] = item
                        st.rerun()

                    col_edit, col_del = st.columns(2)
                    with col_edit:
                        if st.button("✏️ 이름 변경", key=f"rename_{item_id}", use_container_width=True):
                            st.session_state["renaming_history_id"] = item_id
                            st.rerun()
                    with col_del:
                        if st.button("🗑️ 삭제", key=f"delete_{item_id}", use_container_width=True):
                            history = [h for h in history if h["id"] != item_id]
                            save_history(history)
                            if (
                                st.session_state["selected_history"] is not None
                                and st.session_state["selected_history"].get("id") == item_id
                            ):
                                st.session_state["selected_history"] = None
                            st.rerun()

# ── 탭 ───────────────────────────────────────────────────────────────────────
tab1, tab2, tab3 = st.tabs(["📄 심의의견↔소재 1:1 매칭", "✅ 수정 반영 검수", "📏 매체 글자수 체크"])

# ════════════════════════════════════════════════════════════════════════════
# 탭 1: 심의의견 ↔ 소재 1:1 매칭
# ════════════════════════════════════════════════════════════════════════════
with tab1:
    st.subheader("심의의견 ↔ 소재 1:1 매칭")
    st.caption("심의 결과 메일 + 소재 파일을 올리면 각 의견과 소재 위치를 1:1로 매칭해드려요.")

    job_name_input = st.text_input("작업명을 입력하세요", placeholder="예: 삼성화재 암보험 5월 배너 심의")

    col1, col2 = st.columns(2)
    with col1:
        email_text = st.text_area(
            "심의 결과 메일 붙여넣기",
            height=250,
            placeholder="광고주로부터 받은 심의 결과 메일을 전체 복사해서 붙여넣으세요..."
        )
    with col2:
        uploaded_file = st.file_uploader("소재 파일 업로드 (PDF 또는 PPTX)", type=["pdf", "pptx"], key="tab1_file")
        use_ocr = st.toggle("🔍 이미지 OCR 켜기 (느리지만 이미지 텍스트도 읽어요)", value=False)
        if use_ocr:
            st.info("OCR 모드: 이미지 텍스트도 인식합니다. 시간이 걸릴 수 있어요.")

    if st.button("심의의견↔소재 1:1 매칭 시작", type="primary", key="btn_match"):
        if not email_text.strip():
            st.error("심의 결과 메일을 입력해주세요.")
        elif uploaded_file is None:
            st.error("소재 파일을 업로드해주세요.")
        else:
            with st.spinner("파일 분석 중..."):
                file_text = extract_file_text(uploaded_file, use_ocr)

            if not file_text.strip():
                st.warning("⚠️ 파일에서 텍스트를 추출하지 못했습니다. OCR을 켜고 다시 시도해보세요.")
            else:
                with st.spinner("AI가 심의의견과 소재를 매칭하는 중..."):
                    ocr_notice = (
                        "- 이미지에서 추출한 텍스트도 포함되어 있습니다. 이미지 텍스트는 '[슬라이드 N - 이미지 텍스트]' 섹션에 있습니다."
                        if use_ocr else
                        "- OCR이 꺼진 상태입니다. 이미지만 있는 위치는 '🔒 이미지 인식 불가 (OCR 기능을 켜고 다시 시도하세요)'로 표시하세요."
                    )

                    prompt = f"""당신은 보험 광고 심의 전문가입니다.
아래 [심의 결과 메일]과 [소재 파일 텍스트]를 분석해서, 각 심의의견이 소재의 어느 위치에 해당하는지 1:1로 매칭하고 수정 방향을 제안해주세요.

위치 표기 파싱 규칙 (매우 중요):
- "2,7p" → 2페이지 AND 7페이지 (쉼표는 '그리고' 의미, 첫 번째만 아님)
- "2-7p" → 2페이지부터 7페이지
- "_1,6배너" → 배너1 AND 배너6
- "2,7p_1,6배너" → 2페이지와 7페이지의 배너1과 배너6 모두
{ocr_notice}

각 심의의견마다 아래 형식으로 출력하세요:

---
**[의견 N]**
📋 **의견 원문**: (메일에서 추출한 심의의견 원문 그대로)
📄 **페이지 번호**: (해당 페이지 번호, 여러 개면 전부 나열)
🏷️ **배너 번호**: (해당 배너 번호, 없으면 '해당 없음')
❌ **현재 문구**: (소재에서 찾은 현재 문구 또는 이미지 설명)
✅ **수정 문구**: (구체적인 수정 제안)
💡 **수정 이유**: (왜 수정해야 하는지)
---

[심의 결과 메일]
{email_text}

[소재 파일 텍스트]
{file_text}"""

                    try:
                        response = client.chat.completions.create(
                            model="llama-3.3-70b-versatile",
                            messages=[{"role": "user", "content": prompt}],
                            temperature=0.3,
                            max_tokens=4000
                        )
                        result_text = response.choices[0].message.content

                        st.success("매칭 완료!")
                        st.markdown("---")
                        st.markdown(result_text)
                        st.markdown("---")
                        st.download_button(
                            label="📥 결과 저장",
                            data=result_text,
                            file_name="심의의견_매칭결과.txt",
                            mime="text/plain"
                        )

                        import uuid
                        new_item = {
                            "id": str(uuid.uuid4()),
                            "job_name": job_name_input.strip() or "(이름 없음)",
                            "filename": uploaded_file.name,
                            "timestamp": now_kst(),
                            "email_text": email_text,
                            "result_text": result_text,
                        }
                        history = load_history()
                        history.insert(0, new_item)
                        save_history(history)
                        st.rerun()

                    except Exception as e:
                        st.error(f"오류가 발생했습니다: {str(e)}")

# ════════════════════════════════════════════════════════════════════════════
# 탭 2: 수정 반영 검수
# ════════════════════════════════════════════════════════════════════════════
with tab2:
    st.subheader("수정 반영 검수")
    st.caption("이전 심의의견 작업을 불러오고, 수정된 소재 파일을 올리면 각 의견이 반영됐는지 검수해드려요.")

    selected = st.session_state.get("selected_history")

    history = load_history()
    if not history:
        st.info("아직 작업 히스토리가 없어요. 탭 1에서 먼저 심의의견 매칭을 실행해주세요.")
    else:
        history_options = {f"{i}. {h.get('job_name','(이름 없음)')} ({h.get('timestamp','')})": h for i, h in enumerate(history, 1)}
        selected_label = None
        if selected:
            for label, h in history_options.items():
                if h["id"] == selected["id"]:
                    selected_label = label
                    break

        chosen_label = st.selectbox(
            "검수할 작업 선택",
            options=list(history_options.keys()),
            index=list(history_options.keys()).index(selected_label) if selected_label else 0
        )
        chosen_item = history_options[chosen_label]

        st.markdown("**📋 기존 심의의견 요약**")
        with st.expander("심의 결과 메일 보기"):
            st.text(chosen_item.get("email_text", ""))

        revised_file = st.file_uploader("수정된 소재 파일 업로드 (PDF 또는 PPTX)", type=["pdf", "pptx"], key="tab2_file")
        use_ocr2 = st.toggle("🔍 이미지 OCR 켜기", value=False, key="ocr2")

        if st.button("검수 시작", type="primary", key="btn_verify"):
            if revised_file is None:
                st.error("수정된 소재 파일을 업로드해주세요.")
            else:
                with st.spinner("수정 파일 분석 중..."):
                    revised_text = extract_file_text(revised_file, use_ocr2)

                if not revised_text.strip():
                    st.warning("⚠️ 수정 파일에서 텍스트를 추출하지 못했습니다.")
                else:
                    with st.spinner("AI가 심의의견 반영 여부를 검수하는 중..."):
                        prompt = f"""당신은 보험 광고 심의 검수 전문가입니다.
아래 [원래 심의 결과 메일]과 [수정된 소재 파일 텍스트]를 비교해서, 각 심의의견이 제대로 반영됐는지 검수해주세요.

위치 표기 파싱 규칙:
- "2,7p" → 2페이지 AND 7페이지
- "_1,6배너" → 배너1 AND 배너6
- "2,7p_1,6배너" → 2페이지와 7페이지의 배너1과 배너6 모두

각 심의의견마다 아래 형식으로 출력하세요:

---
**[의견 N]**
📋 **의견 원문**: (원래 심의의견 원문)
📄 **페이지 번호**: (해당 페이지)
🏷️ **배너 번호**: (해당 배너, 없으면 '해당 없음')
🔄 **수정 문구**: (수정 파일에서 찾은 현재 문구)
✅ **확인 내용**: (반영 여부 및 세부 내용)
**판정**: ✅ 반영완료 / ❌ 미반영 / ⚠️ 부분반영
---

마지막에 아래 형식으로 총평을 추가하세요:
## 검수 결과 요약
총 N개 중 ✅ N개 / ❌ N개 / ⚠️ N개

[원래 심의 결과 메일]
{chosen_item.get('email_text', '')}

[수정된 소재 파일 텍스트]
{revised_text}"""

                        try:
                            response = client.chat.completions.create(
                                model="llama-3.3-70b-versatile",
                                messages=[{"role": "user", "content": prompt}],
                                temperature=0.3,
                                max_tokens=4000
                            )
                            verify_result = response.choices[0].message.content

                            st.success("검수 완료!")
                            st.markdown("---")
                            st.markdown(verify_result)
                            st.markdown("---")
                            st.download_button(
                                label="📥 검수 결과 저장",
                                data=verify_result,
                                file_name="검수결과.txt",
                                mime="text/plain"
                            )

                        except Exception as e:
                            st.error(f"오류가 발생했습니다: {str(e)}")

# ════════════════════════════════════════════════════════════════════════════
# 탭 3: 매체 글자수 체크
# ════════════════════════════════════════════════════════════════════════════
with tab3:
    st.subheader("매체 글자수 체크")
    media_settings = load_media()

    col_left, col_right = st.columns([1, 1.6])

    with col_left:
        st.markdown("#### 등록된 설정 목록")
        st.caption("항목을 클릭하면 오른쪽에서 글자수 체크를 할 수 있어요.")

        if not media_settings:
            st.info("아래에서 매체/지면을 먼저 등록해주세요.")
        else:
            for idx, s in enumerate(media_settings):
                is_selected = st.session_state["selected_media_idx"] == idx
                is_editing  = st.session_state["editing_idx"] == idx

                if is_editing:
                    with st.container(border=True):
                        st.markdown("**✏️ 편집 중**")
                        e_media     = st.text_input("매체명", value=s["media"],            key=f"e_media_{idx}")
                        e_placement = st.text_input("지면",   value=s["placement"],        key=f"e_place_{idx}")
                        e_category  = st.text_input("구분",   value=s.get("category",""),  key=f"e_cat_{idx}")
                        e_notes     = st.text_input("비고",   value=s.get("notes",""),     key=f"e_notes_{idx}")
                        e_limit     = st.number_input("글자수", value=s["limit"], min_value=1, key=f"e_limit_{idx}")
                        e_spaces    = st.checkbox("공백 포함", value=s["include_spaces"],  key=f"e_spaces_{idx}")
                        c1, c2 = st.columns(2)
                        with c1:
                            if st.button("💾 저장", key=f"save_{idx}", use_container_width=True):
                                media_settings[idx] = {
                                    "media": e_media, "placement": e_placement,
                                    "category": e_category, "notes": e_notes,
                                    "limit": int(e_limit), "include_spaces": e_spaces,
                                }
                                save_media(media_settings)
                                st.session_state["editing_idx"] = None
                                st.rerun()
                        with c2:
                            if st.button("취소", key=f"cancel_{idx}", use_container_width=True):
                                st.session_state["editing_idx"] = None
                                st.rerun()
                else:
                    notes_text  = f" · {s.get('notes','')}" if s.get('notes') else ""
                    space_label = "공백포함" if s["include_spaces"] else "공백제외"
                    prefix      = "✅ " if is_selected else "　 "
                    label       = f"{prefix}{s['media']} — {s.get('category','')}{notes_text}  [{s['limit']}자/{space_label}]"

                    c_btn, c_edit, c_del = st.columns([5, 1, 1])
                    with c_btn:
                        if st.button(label, key=f"sel_{idx}", use_container_width=True):
                            st.session_state["selected_media_idx"] = idx
                            st.rerun()
                    with c_edit:
                        if st.button("✏️", key=f"edit_{idx}", help="편집"):
                            st.session_state["editing_idx"] = idx
                            st.session_state["selected_media_idx"] = idx
                            st.rerun()
                    with c_del:
                        if st.button("🗑", key=f"del_{idx}", help="삭제"):
                            media_settings.pop(idx)
                            save_media(media_settings)
                            st.session_state["selected_media_idx"] = max(0, min(
                                st.session_state["selected_media_idx"], len(media_settings) - 1))
                            st.rerun()

        st.markdown("---")
        with st.expander("➕ 신규 매체/지면 등록"):
            new_media     = st.text_input("매체명",      placeholder="예: 카카오",     key="nm")
            new_placement = st.text_input("지면",        placeholder="예: 배너",       key="np")
            new_category  = st.text_input("구분",        placeholder="예: 광고문구 1", key="nc")
            new_notes     = st.text_input("비고",        placeholder="예: 윗배너",     key="nn")
            new_limit     = st.number_input("글자수 제한", min_value=1, value=20,      key="nl")
            new_spaces    = st.checkbox("공백 포함 글자수", value=True,                key="ns")

            if st.button("등록", type="primary", use_container_width=True):
                if new_media and new_placement and new_category:
                    media_settings.append({
                        "media": new_media, "placement": new_placement,
                        "category": new_category, "notes": new_notes,
                        "limit": int(new_limit), "include_spaces": new_spaces,
                    })
                    save_media(media_settings)
                    st.session_state["selected_media_idx"] = len(media_settings) - 1
                    st.success("등록 완료!")
                    st.rerun()
                else:
                    st.error("매체명, 지면, 구분은 필수입니다.")

    with col_right:
        st.markdown("#### 글자수 체크")
        if not media_settings:
            st.info("왼쪽에서 매체/지면을 먼저 등록해주세요.")
        else:
            idx     = min(st.session_state.get("selected_media_idx") or 0, len(media_settings) - 1)
            setting = media_settings[idx]
            notes   = f" · {setting.get('notes','')}" if setting.get('notes') else ""
            st.info(
                f"**{setting['media']}** — {setting['placement']} — {setting.get('category','')}{notes}  \n"
                f"제한: **{setting['limit']}자** ({'공백 포함' if setting['include_spaces'] else '공백 제외'})"
            )
            check_text = st.text_area("문구 입력", height=220,
                placeholder="글자수를 체크할 문구를 붙여넣으세요...")

            if check_text:
                count = len(check_text) if setting["include_spaces"] else len(check_text.replace(" ", ""))
                limit = setting["limit"]
                space_note = "공백 포함" if setting["include_spaces"] else "공백 제외"
                st.markdown("---")
                if count <= limit:
                    st.success(f"✅ 통과!  **{count}자** / {limit}자 ({space_note})")
                    st.progress(count / limit)
                else:
                    over = count - limit
                    st.error(f"❌ 초과!  **{count}자** / {limit}자 — {over}자 초과 ({space_note})")
                    st.progress(min(count / limit, 1.0))
                st.caption(f"공백 포함: {len(check_text)}자  ·  공백 제외: {len(check_text.replace(' ',''))}자")
